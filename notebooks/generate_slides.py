"""
notebooks/generate_slides.py
Genera reporte/presentacion.pptx  (14 diapositivas, 10–15 min, sinodales).
"""
import sys
from pathlib import Path

_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSAT

# ── Paleta ────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x35, 0x58)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xB2, 0x18, 0x2B)
ORANGE = RGBColor(0xE0, 0x82, 0x14)
BLUE_C = RGBColor(0x43, 0x93, 0xC3)
BLUE_E = RGBColor(0xD6, 0x60, 0x4D)
GREEN  = RGBColor(0x1B, 0x78, 0x37)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)

FIGURAS = _ROOT / "outputs" / "figuras"
OUT     = _ROOT / "reporte" / "presentacion.pptx"

SW = Inches(13.33)   # slide width
SH = Inches(7.5)     # slide height
M  = Inches(0.35)    # margen lateral
HH = Inches(1.05)    # altura del header strip


prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]  # layout completamente en blanco


# ─────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────

def solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color=LGRAY, border=None):
    s = slide.shapes.add_shape(MSAT.RECTANGLE, l, t, w, h)
    solid(s, color)
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.2)
    return s


def txt(slide, text, l, t, w, h,
        size=16, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return bx


def bullets(slide, items, l, t, w, h, size=14,
            text_color=RGBColor(0x22, 0x22, 0x22), spacing=5):
    """items = [(text, indent_level)]  level=0 → bold header"""
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size if level > 0 else size + 1)
        run.font.bold  = (level == 0)
        run.font.color.rgb = text_color


def img(slide, fname, l, t, w=None, h=None):
    path = FIGURAS / fname
    if not path.exists():
        print(f"  ⚠ Figura no encontrada: {fname}")
        return None
    return slide.shapes.add_picture(str(path), l, t, width=w, height=h)


def header_band(slide, title, subtitle=None):
    """Barra azul oscuro fija en la parte superior."""
    rect(slide, 0, 0, SW, HH, NAVY)
    rect(slide, 0, HH - Inches(0.04), SW, Inches(0.04), ORANGE)
    txt(slide, title, M, Inches(0.1), SW - 2*M, Inches(0.72),
        size=22, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, M, Inches(0.76), SW - 2*M, Inches(0.24),
            size=12, color=RGBColor(0xAA, 0xCC, 0xFF))


def stat_box(slide, number, label, l, t, w=Inches(2.9), h=Inches(1.3),
             num_color=ORANGE, bg=NAVY):
    rect(slide, l, t, w, h, bg, border=ORANGE)
    txt(slide, number, l, t + Inches(0.05), w, Inches(0.65),
        size=30, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t + Inches(0.65), w, Inches(0.6),
        size=11, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)


def mk_table(slide, headers, rows, l, t, w, h,
             col_ws=None, hdr_color=NAVY, font_sz=12, alt_row=True):
    cols  = len(headers)
    nrows = len(rows) + 1
    tbl   = slide.shapes.add_table(nrows, cols, l, t, w, h).table
    if col_ws:
        for i, cw in enumerate(col_ws):
            tbl.columns[i].width = cw

    def _cell(cell, text, hdr=False, alt=False):
        cell.text = str(text)
        p  = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_sz)
        run.font.bold = hdr
        run.font.color.rgb = WHITE if hdr else RGBColor(0x1A, 0x1A, 0x1A)
        if hdr:
            cell.fill.solid(); cell.fill.fore_color.rgb = hdr_color
        elif alt:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEB, 0xF2, 0xFA)

    for ci, h_txt in enumerate(headers):
        _cell(tbl.cell(0, ci), h_txt, hdr=True)
    for ri, row in enumerate(rows):
        alt = alt_row and (ri % 2 == 1)
        for ci, val in enumerate(row):
            _cell(tbl.cell(ri + 1, ci), val, alt=alt)
    return tbl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, NAVY)
rect(sl, 0, SH - Inches(0.07), SW, Inches(0.07), ORANGE)
rect(sl, 0, Inches(2.78), SW, Inches(0.04), ORANGE)

txt(sl, "Persistencia Homológica para Detectar",
    M, Inches(0.55), SW - 2*M, Inches(0.85), size=35, bold=True, color=WHITE)
txt(sl, "y Resolver Desiertos de Salud",
    M, Inches(1.35), SW - 2*M, Inches(0.85), size=35, bold=True, color=WHITE)
txt(sl, "Ciudad de México  y  Estado de México",
    M, Inches(2.2), SW - 2*M, Inches(0.55),
    size=22, color=RGBColor(0xAA, 0xCC, 0xFF))

txt(sl, "Mario Carlos Gaitan Reyna",
    M, Inches(3.05), SW - 2*M, Inches(0.45), size=18, bold=True, color=WHITE)
txt(sl, "Datos:  DENUE 2023 (INEGI)  ·  Censo de Población y Vivienda 2020 (INEGI)  ·  OpenStreetMap",
    M, Inches(3.5), SW - 2*M, Inches(0.38),
    size=13, color=RGBColor(0x88, 0xAA, 0xCC))
txt(sl, "Zona Metropolitana del Valle de México",
    M, Inches(3.85), SW - 2*M, Inches(0.38),
    size=13, color=RGBColor(0x88, 0xAA, 0xCC))

rect(sl, M, Inches(4.7), SW - 2*M, Inches(1.1), RGBColor(0x0E, 0x1F, 0x38), border=ORANGE)
txt(sl, ("270,000 personas sin seguro médico atrapadas en desiertos de salud topológicos  ·  "
         "46 huecos críticos CDMX + 77 EDOMEX  ·  K=7 clínicas óptimas para CDMX"),
    M + Inches(0.2), Inches(4.85), SW - 2*M - Inches(0.4), Inches(0.85),
    size=14, color=ORANGE, wrap=True)
print("✓ Slide 1 — Portada")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — El Problema
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "El Problema: Desiertos de Salud Estructurales")

# Bullets izquierda
bullets(sl, [
    ("El problema no es escasez — es distribución", 0),
    ("CDMX tiene 21,585 clínicas; EDOMEX tiene 30,787", 1),
    ("A primera vista no debería haber escasez", 1),
    ("", 0),
    ("Las métricas tradicionales son ciegas", 0),
    ("Densidad de clínicas por km² ✗", 1),
    ("Radio de cobertura de 500 m ✗", 1),
    ("Establecimientos por municipio ✗", 1),
    ("", 0),
    ("Un desierto estructural es invisible para ellas:", 0),
    ("zona rodeada de clínicas por fuera, vacía por dentro", 1),
    ("La persona de adentro puede estar a 2 km de la más cercana", 1),
    ("aunque el mapa de densidad de su colonia se vea 'verde'", 1),
], M, HH + Inches(0.2), Inches(6.6), SH - HH - Inches(1.4), size=14)

# Stats derecha
stat_box(sl, "21,585", "clínicas CDMX",
         SW - Inches(6.4), HH + Inches(0.25), Inches(2.9), Inches(1.3))
stat_box(sl, "30,787", "clínicas EDOMEX",
         SW - Inches(3.2), HH + Inches(0.25), Inches(2.9), Inches(1.3))
stat_box(sl, "270,000", "personas sin seguro\nen desiertos topológicos",
         SW - Inches(4.85), HH + Inches(1.8), Inches(2.9), Inches(1.4),
         num_color=RED)

# Pregunta central
rect(sl, M, SH - Inches(1.35), SW - 2*M, Inches(1.1), NAVY, border=ORANGE)
txt(sl, "Pregunta central: ¿Dónde están esos desiertos, cuánta población atrapan "
        "y cómo priorizamos las inversiones para resolver los más críticos?",
    M + Inches(0.2), SH - Inches(1.26), SW - 2*M - Inches(0.4), Inches(0.9),
    size=14, color=ORANGE, bold=True, wrap=True)
print("✓ Slide 2 — Problema")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Pipeline
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Pipeline Metodológico — Seis Etapas Integradas")

N      = 6
GAP    = Inches(0.18)
BW     = (SW - 2*M - (N - 1) * GAP) / N
BH     = Inches(1.4)
BY     = HH + Inches(0.3)

stages = [
    ("① DATOS",         "DENUE 2023\nCenso 2020\nOpenStreetMap",             BLUE_C),
    ("② TDA",           "Alpha Complex\nLaguerre (pesos)\npersist. > 200 m", NAVY),
    ("③ CRUCE\nCENSAL", "AGEB 2020\nPSINDER\nÍnd. Marginación",             GREEN),
    ("④ RED VIAL",      "OSMnx 4.5 km/h\nIsócronas 15 min\nTiempos reales", RGBColor(0x76,0x2A,0x83)),
    ("⑤ SCORING\n+MCLP","Score 4 ejes\nFiltro OR\nProgramación lineal",     ORANGE),
    ("⑥ VALIDACIÓN",   "Cierre topológico\nRobustez umbral\n100% estable",   RED),
]

for i, (title, desc, color) in enumerate(stages):
    lx = M + i * (BW + GAP)
    rect(sl, lx, BY, BW, BH, color)
    txt(sl, title, lx, BY + Inches(0.06), BW, Inches(0.5),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, lx, BY + Inches(0.55), BW, Inches(0.82),
        size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < N - 1:
        txt(sl, "▶", M + (i + 1)*(BW + GAP) - GAP,
            BY + Inches(0.55), GAP, Inches(0.4),
            size=13, color=GRAY, align=PP_ALIGN.CENTER)

# Fila de resultados clave
results = [
    "21,585 / 30,787\nestab. SCIAN 62",
    "118 CDMX huecos\n241 EDOMEX huecos",
    "62 / 180 habitados\n~270k sin seguro",
    "Umbral > 15 min\nRed peatonal real",
    "46 / 77 prioritarios\n39 / 74 clusters",
    "−14.9% / −16.5%\npersistencia total",
]
for i, res in enumerate(results):
    lx = M + i * (BW + GAP)
    txt(sl, res, lx, BY + BH + Inches(0.12), BW, Inches(0.7),
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Resumen final
rect(sl, M, SH - Inches(1.45), SW - 2*M, Inches(1.08), LGRAY, border=NAVY)
txt(sl,
    "Resultado: K=7 clínicas en CDMX cubren 42,687 personas (44.0%)  ·  "
    "K=9 clínicas ubicables en EDOMEX cubren 35,868 personas (19.0%)  ·  "
    "Validado con persistencia homológica y curvas de robustez",
    M + Inches(0.15), SH - Inches(1.37), SW - 2*M - Inches(0.3), Inches(0.9),
    size=13, color=NAVY, wrap=True)
print("✓ Slide 3 — Pipeline")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Datos
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Los Datos: DENUE 2023 y Censo 2020",
            "Dos fuentes complementarias: dónde están las clínicas y a quién afectan los vacíos")

TY = HH + Inches(0.2)
TH = SH - HH - Inches(0.3)
col_w = (SW - 3*M) / 2

# --- Columna izquierda: DENUE ---
rect(sl, M, TY, col_w, TH, RGBColor(0xE8, 0xF2, 0xFA), border=BLUE_C)
txt(sl, "DENUE 2023 — Establecimientos de Salud",
    M + Inches(0.1), TY + Inches(0.1), col_w - Inches(0.2), Inches(0.4),
    size=14, bold=True, color=NAVY)
bullets(sl, [
    ("Sector SCIAN 62: Servicios de salud", 0),
    ("Hospitales, consultorios, farmacias con consultorio", 1),
    ("Clave PER_OCU: personas ocupadas (tamaño)", 0),
    ("Rango texto → valor numérico", 1),
    ("1–2 pers. hasta 251–500 pers.", 1),
    ("¿Por qué importa el tamaño?", 0),
    ("Un hospital de 300 empleados cubre", 1),
    ("mucho más territorio que un consultorio de 2", 1),
    ("Radio de influencia ∝ √(per_ocu) [Laguerre]", 1),
], M + Inches(0.15), TY + Inches(0.55), col_w - Inches(0.3), TH - Inches(0.7),
   size=13, text_color=RGBColor(0x22, 0x22, 0x22))

# --- Columna derecha: Censo ---
rx = M + col_w + M
rect(sl, rx, TY, col_w, TH, RGBColor(0xF2, 0xFB, 0xF2), border=GREEN)
txt(sl, "Censo 2020 — AGEB Urbanas",
    rx + Inches(0.1), TY + Inches(0.1), col_w - Inches(0.2), Inches(0.4),
    size=14, bold=True, color=NAVY)

mk_table(sl,
    ["Variable", "Qué mide", "Uso"],
    [
        ["PSINDER",    "Personas sin derechohabiencia", "Eje 2 y eje 4"],
        ["GRAPROES",   "Grado promedio escolaridad",    "Índice marg."],
        ["P60YMAS",    "Pob. 60 años y más",            "Índice marg."],
        ["POBTOT",     "Población total",               "Normalización"],
        ["AGEB",       "Área ~manzanas",                "Cruce espacial"],
    ],
    rx + Inches(0.1), TY + Inches(0.55),
    col_w - Inches(0.2), Inches(2.6),
    col_ws=[Inches(1.3), Inches(2.1), Inches(1.2)],
    font_sz=11)

txt(sl, ("Variable crítica: PSINDER\n"
         "Sin seguro médico = sin alternativa cuando la clínica más cercana está lejos\n"
         "270,000 personas en esa situación exacta"),
    rx + Inches(0.1), TY + Inches(3.25), col_w - Inches(0.2), Inches(1.0),
    size=12, color=RED, italic=True)
print("✓ Slide 4 — Datos")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TDA: Alpha Complex y Persistencia
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Análisis Topológico de Datos: Alpha Complex y Persistencia H₁")

# Izquierda: concepto + fórmulas
LW = Inches(5.6)
bullets(sl, [
    ("¿Qué es un hueco H₁?", 0),
    ("Ciclo cerrado de clínicas que rodea un espacio vacío", 1),
    ("Detectado sin importar la forma del contorno", 1),
    ("", 0),
    ("Alpha Complex: construir la topología", 0),
    ("Crecer disco de radio α alrededor de cada clínica", 1),
    ("Aristas cuando dos discos se tocan", 1),
    ("Triángulos cuando tres se tocan", 1),
    ("", 0),
    ("Persistencia = radio_muerte − radio_nacimiento", 0),
    ("Mide qué tan grande es el vacío (en metros)", 1),
    ("Alta persistencia → vacío estructural robusto", 1),
    ("Baja persistencia → ruido entre consultorios", 1),
], M, HH + Inches(0.2), LW, SH - HH - Inches(0.3), size=14)

# Derecha: diagrama de persistencia
img(sl, "persistencia_CDMX.png",
    M + LW + Inches(0.2), HH + Inches(0.15),
    w=SW - M - LW - Inches(0.4))
print("✓ Slide 5 — TDA Alpha Complex")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Weighted Alpha (Laguerre)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Weighted Alpha Complex (Laguerre): El Tamaño Importa",
            "Radio de influencia proporcional al tamaño del establecimiento")

LW = Inches(5.5)
bullets(sl, [
    ("El Alpha clásico trata todo igual", 0),
    ("Farmacia de 2 empleados ≡ Hospital de 300", 1),
    ("Genera falsos positivos donde hay hospitales grandes", 1),
    ("", 0),
    ("Solución: radio proporcional a √(per_ocu)", 0),
    ("r_i = k · √(per_ocu_i)  →  Diagrama de Laguerre", 1),
    ("Cada punto del plano pertenece a la clínica", 1),
    ("que mejor puede cubrirlo (dist. − radio)", 1),
], M, HH + Inches(0.2), LW, Inches(2.6), size=14)

mk_table(sl,
    ["Métrica", "Alpha clásico", "Laguerre", "Cambio"],
    [
        ["Huecos CDMX",   "71",   "118",  "+66.2%"],
        ["Huecos EDOMEX", "209",  "241",  "+15.3%"],
        ["Causa",         "Hospitales grandes\n'rellenan' huecos",
                          "Expone vacíos\nque el hospital\nno alcanza", "—"],
    ],
    M, HH + Inches(2.9), LW, Inches(1.8),
    col_ws=[Inches(1.6), Inches(1.5), Inches(1.3), Inches(1.0)],
    font_sz=12)

txt(sl, "CDMX +66%: mayor heterogeneidad (IMSS, ISSSTE, privados + micro-consultorios)",
    M, SH - Inches(0.8), LW, Inches(0.55),
    size=12, italic=True, color=GRAY)

img(sl, "laguerre_comparacion.png",
    M + LW + Inches(0.2), HH + Inches(0.15),
    w=SW - M - LW - Inches(0.4))
print("✓ Slide 6 — Laguerre")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Calibración min_persistencia = 200 m
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Calibración: min_persistencia = 200 m y Robustez",
            "Punto de equilibrio entre señal estructural y ruido topológico")

LW = Inches(5.5)
mk_table(sl,
    ["Umbral", "Huecos CDMX", "Huecos EDOMEX", "Interpretación"],
    [
        ["100 m", "482", "742",  "Micro-brechas entre consultorios adyacentes"],
        ["200 m ✓", "118", "241", "Señal estructural: vacíos ≥ 1 cuadra larga"],
        ["400 m", "17",  "91",   "Solo los más severos; pierde huecos medianos"],
        ["600 m", "7",   "58",   "Únicamente vacíos de escala de colonia"],
    ],
    M, HH + Inches(0.25), LW, Inches(2.2),
    col_ws=[Inches(1.0), Inches(1.1), Inches(1.3), Inches(2.0)],
    font_sz=12)

rect(sl, M, HH + Inches(2.55), LW, Inches(0.95),
     RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
txt(sl, ("Prueba de robustez:  "
         "Los 7 huecos de 600 m en CDMX aparecen también a 100, 200 y 400 m (100%).\n"
         "Los 58 de EDOMEX ídem. Los huecos más severos son señales topológicas sólidas."),
    M + Inches(0.1), HH + Inches(2.62), LW - Inches(0.2), Inches(0.82),
    size=12, color=GREEN)

txt(sl, ("200 m = un espacio donde ninguna clínica está a < 200 m.\n"
         "Para el interior del vacío, la distancia real puede ser > 500 m."),
    M, HH + Inches(3.65), LW, Inches(0.7),
    size=12, italic=True, color=GRAY)

img(sl, "robustez_conteo_huecos.png",
    M + LW + Inches(0.2), HH + Inches(0.15),
    w=SW - M - LW - Inches(0.4))
print("✓ Slide 7 — Calibración")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Resultados topológicos + Impacto demográfico
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Resultados Topológicos e Impacto Demográfico",
            "Del algoritmo a las personas: ¿quién vive dentro de los vacíos?")

# Stat boxes fila 1
stat_box(sl, "118",  "huecos topológicos\nCDMX",       M,                HH + Inches(0.2),
         bg=NAVY,   num_color=BLUE_C)
stat_box(sl, "241",  "huecos topológicos\nEDOMEX",     M + Inches(3.1),  HH + Inches(0.2),
         bg=NAVY,   num_color=BLUE_E)
stat_box(sl, "62/180", "habitados (filtro Censo\nCDMX / EDOMEX)", M + Inches(6.2), HH + Inches(0.2),
         bg=RGBColor(0x1B,0x47,0x30), num_color=WHITE)
stat_box(sl, "270,000", "personas sin seguro\nen desiertos topológicos",
         M + Inches(9.3), HH + Inches(0.2),
         bg=RED, num_color=WHITE, w=Inches(3.6))

# Texto aclaratorio
txt(sl, ("CDMX: Cluster 1 (poniente) → 7,525 sin seguro, clínica más cercana > 12 min caminando\n"
         "EDOMEX: Cluster 5 (norponiente) → 23,869 sin seguro en 3 huecos; 33.7% de la zona sin derechohabiencia"),
    M, HH + Inches(1.75), SW - 2*M, Inches(0.8), size=13, color=NAVY)

# Tabla resumen por región
mk_table(sl,
    ["Región", "Huecos con pobl.", "Personas sin seguro", "Concentración"],
    [
        ["CDMX",    "62",  "~88,000",  "Alta en zonas periféricas"],
        ["EDOMEX",  "180", "~182,000", "Mayor dispersión geográfica"],
        ["TOTAL",   "242", "~270,000", "≈ población de una ciudad media MX"],
    ],
    M, HH + Inches(2.65), Inches(7.5), Inches(1.7),
    col_ws=[Inches(1.2), Inches(1.6), Inches(1.8), Inches(2.8)],
    font_sz=12)

img(sl, "matriz_riesgo_CDMX.png",
    M + Inches(7.8), HH + Inches(1.5),
    w=Inches(5.0))
print("✓ Slide 8 — Resultados + Impacto")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Accesibilidad Real: Red Vial OSMnx
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Accesibilidad Real: Red Vial Peatonal con OSMnx",
            "Distancia euclidiana × 1.35 subestima barreras reales: barrancas, vías rápidas, zonas industriales")

LW = Inches(5.5)
bullets(sl, [
    ("¿Por qué la distancia euclidiana no basta?", 0),
    ("800 m en línea recta ≠ 800 m caminando", 1),
    ("Vías rápidas sin cruce, barrancas, muros", 1),
    ("", 0),
    ("OSMnx: red vial peatonal de OpenStreetMap", 0),
    ("Velocidad estándar: 4.5 km/h peatonal", 1),
    ("Isócronas = polígonos alcanzables en t minutos", 1),
    ("concave_hull(ratio=0.4) → forma real de la red", 1),
    ("", 0),
    ("Umbral de intervención: 15 minutos", 0),
    ("Respaldado por literatura de salud pública", 1),
    ("Reducción significativa de consulta preventiva", 1),
    ("cuando la barrera supera este umbral", 1),
    ("", 0),
    ("Huecos más críticos:", 0),
    ("CDMX: hasta 35 min (Álvaro Obregón, Xochimilco)", 1),
    ("EDOMEX: más de 40 min (zonas norte y oriente)", 1),
], M, HH + Inches(0.2), LW, SH - HH - Inches(0.3), size=13)

img(sl, "isocronas_roma_norte.png",
    M + LW + Inches(0.2), HH + Inches(0.1),
    w=SW - M - LW - Inches(0.35),
    h=SH - HH - Inches(0.25))
print("✓ Slide 9 — Accesibilidad OSMnx")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Scoring 4 ejes + Filtro OR → 46/77 prioritarios
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Scoring Multi-eje y Selección de Huecos Prioritarios",
            "score = 0.25·r_tiempo + 0.25·r_sin_seguro + 0.25·r_persistencia + 0.25·r_marginación")

LW = Inches(5.5)

mk_table(sl,
    ["Eje", "Variable", "Peso", "Umbral filtro OR"],
    [
        ["① Inaccesibilidad", "Tiempo caminata (min)",          "25%", "> 10 min"],
        ["② Inequidad",       "Personas sin seguro",            "25%", "> percentil 60"],
        ["③ Topológico",      "Persistencia del hueco (m)",     "25%", "> 400 m"],
        ["④ Marginación",     "Índice compuesto Censo 2020",    "25%", "> percentil 60"],
    ],
    M, HH + Inches(0.2), LW, Inches(2.1),
    col_ws=[Inches(1.7), Inches(1.8), Inches(0.7), Inches(1.2)],
    font_sz=12)

bullets(sl, [
    ("Ventaja del percentil: evita que un eje domine", 0),
    ("Filtro OR: basta cumplir 1 de 4 ejes → entra", 0),
    ("Descartados: micro-brechas bajas en los 4 simultáneamente", 0),
    ("DBSCAN (eps=1500 m): huecos cercanos = una decisión de inversión", 0),
], M, HH + Inches(2.4), LW, Inches(1.8), size=13)

# Stats de resultado
stat_box(sl, "46", "huecos prioritarios\nCDMX  →  39 clusters", M,
         SH - Inches(1.6), w=Inches(2.5), h=Inches(1.3), bg=NAVY, num_color=BLUE_C)
stat_box(sl, "77", "huecos prioritarios\nEDOMEX  →  74 clusters",
         M + Inches(2.7), SH - Inches(1.6), w=Inches(2.5), h=Inches(1.3),
         bg=NAVY, num_color=BLUE_E)

img(sl, "clusters_geograficos.png",
    M + LW + Inches(0.2), HH + Inches(0.1),
    w=SW - M - LW - Inches(0.35))
print("✓ Slide 10 — Scoring y Clusters")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MCLP: formulación + CDMX K=7 curva y mapa
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "CDMX — Optimización MCLP: ¿Cuántas Clínicas y Dónde?",
            "Maximum Coverage Location Problem · demanda = score × personas_sin_seguro · cobertura ≤ 15 min")

LW = Inches(5.4)
TY = HH + Inches(0.15)

txt(sl, "MCLP: max Σ (score_i × pob_sin_salud_i) · y_i  s.t.  Σ x_j = K,  y_i ≤ Σ C_ij·x_j",
    M, TY, LW, Inches(0.5), size=12, bold=True, color=NAVY, italic=True)

mk_table(sl,
    ["K", "% cubierto", "Incremento"],
    [
        ["5",   "33.0%", "—"],
        ["6",   "38.6%", "+5.6 pp"],
        ["7 ✓", "44.0%", "+5.3 pp"],
        ["8",   "49.1%", "+5.1 pp"],
        ["9",   "53.1%", "+4.0 pp ↓"],
        ["10",  "56.1%", "+3.0 pp ↓"],
    ],
    M, TY + Inches(0.6), LW, Inches(2.85),
    col_ws=[Inches(0.65), Inches(1.3), Inches(1.35)],
    font_sz=13)

rect(sl, M, TY + Inches(3.55), LW, Inches(1.0),
     RGBColor(0xE8, 0xF2, 0xFA), border=BLUE_C)
txt(sl, ("K=7: último punto con >5 pp por clínica.\n"
         "K=8→9 muestra la caída más pronunciada (−1.1 pp):\n"
         "los huecos de mayor densidad ya están cubiertos."),
    M + Inches(0.1), TY + Inches(3.62), LW - Inches(0.2), Inches(0.9),
    size=13, color=NAVY)

img(sl, "fase3_CDMX_k7_mapa.png",
    M + LW + Inches(0.15), HH + Inches(0.1),
    w=SW - M - LW - Inches(0.3))
print("✓ Slide 11 — MCLP CDMX K=7")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CDMX K=7: análisis de los 4 ejes
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "CDMX K=7: Análisis de los 4 Ejes — ¿Qué Cubre y Qué Persiste?",
            "11 huecos cubiertos · 35 persistentes · 42,687 personas sin seguro cubiertas (44.0%)")

IMG_H = Inches(3.9)
img(sl, "fase3_CDMX_k7_ejes.png",
    M, HH + Inches(0.1), w=SW - 2*M, h=IMG_H)

BY = HH + IMG_H + Inches(0.25)
mk_table(sl,
    ["Eje", "Cubiertos (11)", "Persistentes (35)", "Lectura"],
    [
        ["Tiempo (min)",       "7.89",  "9.27",  "MCLP cubre los más accesibles con más población"],
        ["Sin seguro (pers.)", "3,881", "1,555", "Priorización correcta por inequidad (×2.5)"],
        ["Persistencia (m)",  "296",   "364",   "Vacíos más grandes quedan fuera (zonas periféricas)"],
        ["Marginación (IM)",  "0.35",  "0.29",  "El 4.° eje redirige a zonas más vulnerables"],
    ],
    M, BY, SW - 2*M, Inches(1.95),
    col_ws=[Inches(1.8), Inches(1.6), Inches(1.75), Inches(5.5)],
    font_sz=11)
print("✓ Slide 12 — CDMX K=7 Ejes")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — EDOMEX: estrategia subregional K=12
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "EDOMEX — Por Qué una Estrategia Distinta: K=12 Subregional",
            "Curva lineal sin codo → la solución no es un K global sino dividir el problema")

LW = Inches(5.4)
TY = HH + Inches(0.15)

bullets(sl, [
    ("La curva K=5→10 de EDOMEX es casi lineal:", 0),
    ("~+2 pp por clínica; no hay codo ni punto óptimo único", 1),
    ("74 de 77 clusters son singletons → alta dispersión", 1),
    ("ZMVM, Norte, Poniente y Oriente no se comunican a 15 min", 1),
    ("", 0),
    ("Solución: K=3 por zona (equidad geográfica)", 0),
    ("K=3 uniforme corrige que ZMVM ya tiene más infraestructura", 1),
    ("Norte y Poniente tienen huecos más persistentes estructuralmente", 1),
], M, TY, LW, Inches(2.4), size=13)

mk_table(sl,
    ["Zona", "Huecos", "K", "Sin seguro cub.", "% local"],
    [
        ["Norte",       "8",  "3", "1,182",  "13.9%"],
        ["Poniente",    "22", "3", "10,243", "23.0%"],
        ["ZMVM-Centro", "41", "3", "24,443", "19.7%"],
        ["Oriente",     "6",  "3", "Sin candidatos válidos", "—"],
        ["TOTAL",       "77", "9 ubicables", "35,868", "19.0%"],
    ],
    M, TY + Inches(2.5), LW, Inches(2.35),
    col_ws=[Inches(1.3), Inches(0.75), Inches(0.75), Inches(1.55), Inches(0.95)],
    font_sz=12)

txt(sl, "Oriente: 6 huecos, 12,270 sin seguro. Requiere umbral extendido o clínicas móviles.",
    M, SH - Inches(0.55), LW, Inches(0.4),
    size=11, italic=True, color=RED)

img(sl, "fase3_EDOMEX_subregional_mapa.png",
    M + LW + Inches(0.15), HH + Inches(0.1),
    w=SW - M - LW - Inches(0.3))
print("✓ Slide 13 — EDOMEX Subregional")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Validación + Conclusiones
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_band(sl, "Validación Topológica y Conclusiones")

LW = Inches(5.8)
TY = HH + Inches(0.2)

# Validación topológica
txt(sl, "Cierre topológico: ¿las clínicas nuevas realmente cierran los huecos?",
    M, TY, LW, Inches(0.4), size=14, bold=True, color=NAVY)

mk_table(sl,
    ["Estado", "Condición", "CDMX K=7", "EDOMEX K=12"],
    [
        ["Cerrado",    "d < pers_m",       "3 / 46  (7%)",  "6 / 77  (8%)"],
        ["Parcial",    "pers_m ≤ d < 2·r", "4 / 46  (9%)",  "2 / 77  (3%)"],
        ["Persistente","d ≥ 2·pers_m",     "39 / 46 (85%)", "69 / 77 (90%)"],
        ["Δ persist.",  "reducción total",  "−14.9%",        "−16.5%"],
    ],
    M, TY + Inches(0.45), LW, Inches(2.2),
    col_ws=[Inches(1.3), Inches(1.55), Inches(1.35), Inches(1.45)],
    font_sz=12)

rect(sl, M, TY + Inches(2.75), LW, Inches(0.8),
     RGBColor(0xFF, 0xF8, 0xE1), border=ORANGE)
txt(sl, ("MCLP y cierre topológico miden cosas distintas: MCLP optimiza acceso\n"
         "poblacional (≤15 min caminando). Cierre topológico requiere d < pers_m del centroide."),
    M + Inches(0.1), TY + Inches(2.8), LW - Inches(0.2), Inches(0.7),
    size=11, italic=True, color=RGBColor(0x80, 0x40, 0x00))

# Conclusiones
txt(sl, "Conclusiones clave",
    M, TY + Inches(3.65), LW, Inches(0.35), size=14, bold=True, color=NAVY)
bullets(sl, [
    ("El problema es distribución, no escasez", 0),
    ("21k / 30k clínicas con 118 / 241 vacíos estructurales", 1),
    ("La topología revela lo que la densidad no ve", 0),
    ("+66% huecos vs Alpha clásico en CDMX con Laguerre", 1),
    ("El 4.° eje (marginación) cambia quién recibe inversión", 0),
    ("IM cubiertos 0.35 > 0.29 persistentes → más justo", 1),
], M, TY + Inches(4.05), LW, Inches(1.8), size=13)

img(sl, "validacion_mapa.png",
    M + LW + Inches(0.15), HH + Inches(0.1),
    w=SW - M - LW - Inches(0.3))
print("✓ Slide 14 — Validación + Conclusiones")


# ═══════════════════════════════════════════════════════════════════════════
# Guardar
# ═══════════════════════════════════════════════════════════════════════════
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"\n✓ Presentación guardada: {OUT}")
print(f"  {len(prs.slides)} diapositivas")
